"""
Document parsing for the Context Catalog.

Converts uploaded documents (docx / pptx / pdf) and raw markdown/text into a
single markdown/plain-text string used for retrieval. Parser libraries are
imported lazily so the rest of the app keeps working if an optional dependency
is missing (the upload endpoint surfaces a clear error instead of crashing on
import).
"""
import logging
import os
from typing import Dict, Callable

logger = logging.getLogger(__name__)


class DocumentParseError(Exception):
    """Raised when a document cannot be parsed into text."""


# extension -> logical doc_type stored on the document row
SUPPORTED_UPLOAD_EXTENSIONS: Dict[str, str] = {
    ".md": "markdown",
    ".markdown": "markdown",
    ".txt": "markdown",
    ".docx": "docx",
    ".pptx": "pptx",
    ".pdf": "pdf",
}


def detect_doc_type(filename: str) -> str:
    """Map a filename to a logical doc_type, or raise if unsupported."""
    ext = os.path.splitext(filename or "")[1].lower()
    if ext not in SUPPORTED_UPLOAD_EXTENSIONS:
        supported = ", ".join(sorted(SUPPORTED_UPLOAD_EXTENSIONS))
        raise DocumentParseError(
            f"Unsupported file type '{ext or filename}'. Supported types: {supported}"
        )
    return SUPPORTED_UPLOAD_EXTENSIONS[ext]


def _parse_text(content: bytes) -> str:
    try:
        return content.decode("utf-8")
    except UnicodeDecodeError:
        return content.decode("latin-1", errors="replace")


def _parse_docx(content: bytes) -> str:
    try:
        from docx import Document  # python-docx
    except ImportError as e:  # pragma: no cover - depends on optional dep
        raise DocumentParseError(
            "python-docx is not installed; cannot parse .docx files."
        ) from e
    from io import BytesIO

    document = Document(BytesIO(content))
    lines = []
    for para in document.paragraphs:
        text = (para.text or "").strip()
        if not text:
            continue
        style = (para.style.name or "").lower() if para.style else ""
        if style.startswith("heading"):
            # Best-effort: map Word heading levels to markdown headings.
            level = "".join(ch for ch in style if ch.isdigit())
            prefix = "#" * (int(level) if level.isdigit() else 2)
            lines.append(f"{prefix} {text}")
        else:
            lines.append(text)

    # Append simple table content as pipe rows so it remains searchable.
    for table in document.tables:
        for row in table.rows:
            cells = [(cell.text or "").strip() for cell in row.cells]
            if any(cells):
                lines.append(" | ".join(cells))

    return "\n\n".join(lines).strip()


def _parse_pptx(content: bytes) -> str:
    try:
        from pptx import Presentation  # python-pptx
    except ImportError as e:  # pragma: no cover - depends on optional dep
        raise DocumentParseError(
            "python-pptx is not installed; cannot parse .pptx files."
        ) from e
    from io import BytesIO

    prs = Presentation(BytesIO(content))
    blocks = []
    for idx, slide in enumerate(prs.slides, start=1):
        slide_lines = [f"## Slide {idx}"]
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for para in shape.text_frame.paragraphs:
                text = "".join(run.text for run in para.runs).strip()
                if text:
                    slide_lines.append(text)
        if len(slide_lines) > 1:
            blocks.append("\n".join(slide_lines))
    return "\n\n".join(blocks).strip()


def _parse_pdf(content: bytes) -> str:
    try:
        from pypdf import PdfReader
    except ImportError as e:  # pragma: no cover - depends on optional dep
        raise DocumentParseError(
            "pypdf is not installed; cannot parse .pdf files."
        ) from e
    from io import BytesIO

    reader = PdfReader(BytesIO(content))
    pages = []
    for page in reader.pages:
        try:
            text = (page.extract_text() or "").strip()
        except Exception as e:  # noqa: BLE001 - per-page extraction is best-effort
            logger.warning("Failed to extract text from a PDF page: %s", e)
            text = ""
        if text:
            pages.append(text)
    return "\n\n".join(pages).strip()


_PARSERS: Dict[str, Callable[[bytes], str]] = {
    "markdown": _parse_text,
    "docx": _parse_docx,
    "pptx": _parse_pptx,
    "pdf": _parse_pdf,
}


def parse_document(filename: str, content: bytes) -> str:
    """Extract markdown/plain text from an uploaded document.

    Args:
        filename: Original filename (used to detect the type).
        content: Raw file bytes.

    Returns:
        Extracted markdown/plain text.

    Raises:
        DocumentParseError: if the type is unsupported, a parser dependency is
            missing, or no text could be extracted.
    """
    doc_type = detect_doc_type(filename)
    parser = _PARSERS[doc_type]
    try:
        text = parser(content)
    except DocumentParseError:
        raise
    except Exception as e:  # noqa: BLE001 - normalize parser failures
        raise DocumentParseError(f"Failed to parse '{filename}': {e}") from e

    if not text or not text.strip():
        raise DocumentParseError(
            f"No extractable text found in '{filename}'. "
            "Scanned/image-only documents are not supported."
        )
    return text.strip()
